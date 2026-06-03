"""PPT 导出工具：模板化 16:9 表格输出。

依赖：python-pptx

设计要点：
- 表头采用品牌色 + 白字 + 居中加粗
- 数据行斑马纹（浅灰底）+ 边框 + 左对齐 + 自动换行
- 支持父子分组表头：通过将相邻列的父表头合并并设置同样的填充色实现
- 标题区域含主标题 + 副标题（导出时间）
"""
import io
import json
from datetime import datetime
from typing import Iterable, List, Optional, Sequence


def checklist_to_text(val: str) -> str:
    """将清单字段值（JSON 或旧纯文本）转为可读字符串，用于 PPT/导出。
    格式：每行前缀 ✓（已完成）或 ·（未完成）。
    """
    if not val:
        return ""
    try:
        items = json.loads(val)
        if isinstance(items, list):
            lines = []
            for item in items:
                text = str(item.get("text", "")).strip()
                if text:
                    lines.append(("✓ " if item.get("done") else "· ") + text)
            return "\n".join(lines)
    except (ValueError, TypeError, AttributeError):
        pass
    return val  # 旧纯文本，原样返回

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt


_SLIDE_W = Emu(12192000)  # 16:9 默认 13.333"
_SLIDE_H = Emu(6858000)

# 华为风格：标志性红 #C7000B，深红压顶 + 浅红斑马
_BRAND = RGBColor(0xC7, 0x00, 0x0B)        # 华为红
_BRAND_DARK = RGBColor(0x9E, 0x00, 0x09)   # 深红（横幅）
_ZEBRA = RGBColor(0xFB, 0xF2, 0xF2)        # 极浅红灰斑马
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_TEXT = RGBColor(0x26, 0x26, 0x26)
_BORDER = RGBColor(0xE3, 0xC9, 0xCB)       # 浅红灰边框
_SUBTITLE = RGBColor(0xF2, 0xD0, 0xD2)     # 横幅副标题浅红

# 中文 / 西文字体（华为优先 HarmonyOS Sans，回退微软雅黑）
_FONT_LATIN = "HarmonyOS Sans SC"
_FONT_EA = "微软雅黑"

# 进展状态着色（领域/产品需求 slide 自动染色，提升可读性）
_STATUS_COLORS = {
    "已完成": RGBColor(0x2E, 0x7D, 0x32),
    "进行中": RGBColor(0x15, 0x65, 0xC0),
    "已延期": RGBColor(0xC6, 0x28, 0x28),
    "已变更": RGBColor(0xB9, 0x6A, 0x00),
    "未开始": RGBColor(0x90, 0x93, 0x99),
    "不涉及": RGBColor(0xB0, 0xB3, 0xB8),
}


def _new_pres() -> Presentation:
    pres = Presentation()
    pres.slide_width = _SLIDE_W
    pres.slide_height = _SLIDE_H
    return pres


def _add_banner(slide, title: str, subtitle: str):
    """顶部品牌色横幅：主标题 + 副标题。"""
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), _SLIDE_W, Inches(0.8))
    banner.line.fill.background()
    banner.fill.solid()
    banner.fill.fore_color.rgb = _BRAND_DARK
    banner.shadow.inherit = False

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(10), Inches(0.45))
    tf = title_box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _apply_run_font(p.runs[0], 22, True, _WHITE)

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.5), Inches(10), Inches(0.3))
    tf2 = sub_box.text_frame
    tf2.margin_left = tf2.margin_right = 0
    tf2.text = subtitle
    p2 = tf2.paragraphs[0]
    _apply_run_font(p2.runs[0], 11, False, _SUBTITLE)


def _fit_font_size(rows_count: int) -> int:
    if rows_count <= 6:
        return 11
    if rows_count <= 12:
        return 10
    if rows_count <= 20:
        return 9
    return 8


def _rgb_to_hex(color: RGBColor) -> str:
    return str(color)  # python-pptx RGBColor 的 __str__ 返回 6 位大写十六进制


def _apply_run_font(run, size: int, bold: bool, color: RGBColor):
    """统一设置 run 的字号/粗细/颜色，并补齐东亚字体（中文不走默认衬线）。"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = _FONT_LATIN
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", _FONT_EA)


def _set_cell_margins(cell, lr=5, tb=2):
    cell.margin_left = Pt(lr)
    cell.margin_right = Pt(lr)
    cell.margin_top = Pt(tb)
    cell.margin_bottom = Pt(tb)


def _clear_table_style(table):
    """清掉 python-pptx 默认套用的主题表样式（带蓝色条纹），
    换成"无样式"，让我们手动设置的填充/边框完全生效。"""
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn("a:tblPr"))
    for attr in ("firstRow", "lastRow", "firstCol", "lastCol", "bandRow", "bandCol"):
        tblPr.set(attr, "0")
    style_id = tblPr.find(qn("a:tableStyleId"))
    if style_id is None:
        style_id = etree.SubElement(tblPr, qn("a:tableStyleId"))
    # "No Style, No Grid"
    style_id.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _set_cell_border(cell, color: RGBColor = _BORDER):
    """给单元格四边加细边框。python-pptx 没有现成 API，用 lxml 操作。"""
    tc_pr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        # 先删除已有，避免重复
        existing = tc_pr.findall(qn(tag))
        for e in existing:
            tc_pr.remove(e)
    hex_val = _rgb_to_hex(color)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = etree.SubElement(tc_pr, qn(tag))
        ln.set("w", "6350")  # 0.5pt
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        solid = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(solid, qn("a:srgbClr"))
        srgb.set("val", hex_val)


def _style_header_cell(cell, text: str, font_size: int):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = _BRAND
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_cell_margins(cell)
    for para in cell.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for r in para.runs:
            _apply_run_font(r, font_size + 1, True, _WHITE)
    _set_cell_border(cell, _WHITE)


def _style_data_cell(cell, value, font_size: int, zebra: bool, center: bool = False):
    text = "" if value is None else str(value)
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = _ZEBRA if zebra else _WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_cell_margins(cell)

    # 进展状态自动染色 + 居中加粗；短列居中
    stripped = text.strip()
    status_color = _STATUS_COLORS.get(stripped)
    is_status = status_color is not None
    color = status_color if is_status else _TEXT

    for para in cell.text_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER if (center or is_status) else PP_ALIGN.LEFT
        for r in para.runs:
            _apply_run_font(r, font_size, is_status, color)
    _set_cell_border(cell)


def _merge_header_groups(table, header_row: int, groups: Sequence[tuple]):
    """合并某一表头行的相邻列：groups = [(start_col, end_col, label), ...]"""
    for start, end, label in groups:
        if end > start:
            table.cell(header_row, start).merge(table.cell(header_row, end))
        cell = table.cell(header_row, start)
        cell.text = label


def _add_grouped_table(
    slide,
    parent_headers: Sequence[Optional[tuple]],  # list aligned with cols: (group_id, label) or None
    leaf_headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    col_widths_in: Optional[Sequence[float]] = None,
    center_cols: Optional[set] = None,
):
    """支持「父表头 + 子表头」两行表头的表格。

    parent_headers: 长度与 leaf_headers 相同。每个元素是 (group_id, label) 或 None。
        相邻同 group_id 会被合并；None 表示该列没有父分组，会与子表头单元格垂直合并显示。
    center_cols: 需要居中显示的数据列下标集合（短列，如序号/优先级）。
    """
    center_cols = center_cols or set()
    n_cols = len(leaf_headers)
    n_rows = 2 + len(rows)  # 2 行表头

    left = Inches(0.4)
    top = Inches(0.95)
    width = Inches(12.5)
    height = Inches(5.7)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    _clear_table_style(table)

    if col_widths_in and len(col_widths_in) == n_cols:
        for i, w in enumerate(col_widths_in):
            table.columns[i].width = Inches(w)

    font_size = _fit_font_size(len(rows))

    # ===== 父表头行 =====
    # 先填默认 label（空），稍后合并
    for j in range(n_cols):
        _style_header_cell(table.cell(0, j), "", font_size)

    # 计算分组
    groups: List[tuple] = []  # (start, end, label, group_id)
    j = 0
    while j < n_cols:
        ph = parent_headers[j]
        if ph is None:
            # 此列父行将与子行做垂直合并
            j += 1
            continue
        gid, label = ph
        end = j
        while end + 1 < n_cols and parent_headers[end + 1] is not None and parent_headers[end + 1][0] == gid:
            end += 1
        groups.append((j, end, label, gid))
        j = end + 1

    _merge_header_groups(table, 0, [(s, e, lab) for s, e, lab, _ in groups])

    # 重新设置合并后的父单元格样式（merge 之后需要再次设置文本/样式）
    for s, e, lab, _ in groups:
        _style_header_cell(table.cell(0, s), lab, font_size)

    # ===== 子表头行 =====
    for j, h in enumerate(leaf_headers):
        _style_header_cell(table.cell(1, j), h, font_size)

    # 把「没有父分组」的列做垂直合并（父行 + 子行 -> 同一格显示子标题）
    for j in range(n_cols):
        if parent_headers[j] is None:
            top_cell = table.cell(0, j)
            bot_cell = table.cell(1, j)
            top_cell.merge(bot_cell)
            _style_header_cell(table.cell(0, j), leaf_headers[j], font_size)

    # ===== 数据行 =====
    for i, row in enumerate(rows, start=2):
        zebra = (i - 2) % 2 == 1
        for j, val in enumerate(row):
            _style_data_cell(table.cell(i, j), val, font_size, zebra, center=(j in center_cols))


def build_customer_status_pptx(rows: Iterable) -> io.BytesIO:
    """传入 CustomerStatus ORM 列表，返回 BytesIO。"""
    leaf_headers = [
        "机台编号", "客户", "型号", "当前阶段", "现场版本", "关注度",
        "当前进展", "现场关键事务", "软件类风险和问题", "问题单",
    ]
    parent_headers: List[Optional[tuple]] = [None] * len(leaf_headers)
    col_widths = [0.85, 1.0, 0.9, 0.95, 0.95, 0.7, 2.0, 2.0, 2.0, 1.15]

    data: List[List[str]] = []
    for r in rows:
        data.append([
            r.machine_id or "",
            r.battlefield or "",
            r.model or "",
            r.current_stage or "",
            r.field_version or "",
            ("★" * (r.attention_level or 0)) or "-",
            r.customer_status or "",
            checklist_to_text(r.recent_focus or ""),
            checklist_to_text(r.key_issues or ""),
            getattr(r, "issue_url", "") or "—",
        ])

    pres = _new_pres()
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # 空白
    _add_banner(
        slide,
        "客户面状态总览",
        f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}   ·   共 {len(data)} 条",
    )
    _add_grouped_table(slide, parent_headers, leaf_headers, data, col_widths,
                       center_cols={0, 2, 3, 4, 5, 9})

    buf = io.BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf


def _iteration_title(iteration) -> str:
    title = f"{iteration.year}年{iteration.month}月迭代"
    if iteration.name:
        title += f"  ·  {iteration.name}"
    return title


def _add_domain_slide(pres, iteration, requirements: Iterable):
    """领域需求 slide：基础列 + 「交付进展跟踪」6 子列 + 备注。"""
    leaf_headers = [
        "序号", "需求编号", "需求标题", "责任人", "PL组", "优先级", "计划版本",
        "需求串讲", "反串讲", "STC设计", "编码", "BBIT", "转测澄清",
        "备注",
    ]
    parent_headers: List[Optional[tuple]] = [None] * 7 + [("progress", "交付进展跟踪")] * 6 + [None]
    col_widths = [0.45, 1.0, 2.0, 0.75, 0.7, 0.55, 0.85, 0.82, 0.82, 0.82, 0.82, 0.82, 0.82, 1.4]

    data: List[List[str]] = []
    for r in requirements:
        data.append([
            str(r.seq or 0),
            r.req_no or "",
            r.title or "",
            r.owner or "",
            getattr(r, "owner_group", "") or "",
            r.priority or "",
            r.planned_version or "",
            r.progress_walkthrough or "",
            r.progress_reverse or "",
            r.progress_stc or "",
            r.progress_coding or "",
            r.progress_bbit or "",
            r.progress_clarify or "",
            getattr(r, "remark", "") or "",
        ])

    slide = pres.slides.add_slide(pres.slide_layouts[6])
    title = _iteration_title(iteration) + "  ·  领域需求"
    subtitle = f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}   ·   共 {len(data)} 条领域需求"
    if iteration.owner:
        subtitle += f"   ·   负责人：{iteration.owner}"
    _add_banner(slide, title, subtitle)
    _add_grouped_table(slide, parent_headers, leaf_headers, data, col_widths,
                       center_cols={0, 5, 7, 8, 9, 10, 11, 12})


def _add_product_slides(pres, iteration, product_reqs: Iterable):
    """产品需求 slide：单页字段太多，拆 2 张：基础+特性 / 交付进展跟踪。"""
    rows_basic: List[List[str]] = []
    rows_progress: List[List[str]] = []
    for r in product_reqs:
        rows_basic.append([
            str(r.seq or 0),
            r.req_no or "",
            r.title or "",
            r.planned_version or "",
            r.priority or "",
            r.feature or "",
            r.feature_fo or "",
            r.feature_se or "",
            r.feature_tfo or "",
            r.code_areas or "",
            r.key_risks or "",
        ])
        rows_progress.append([
            str(r.seq or 0),
            r.req_no or "",
            r.title or "",
            r.progress_walkthrough or "",
            r.progress_reverse or "",
            r.progress_domain or "",
            r.progress_coding or "",
            r.progress_joint_debug or "",
            r.progress_clarify or "",
            r.progress_test_result or "",
            r.estimated_loc or "",
            r.actual_loc or "",
            r.actual_effort or "",
        ])

    base_title = _iteration_title(iteration)
    subtitle_common = f"导出时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}   ·   共 {len(rows_basic)} 条产品需求"

    # —— Slide 1：基础信息 + 特性 + 关键风险 ——
    slide1 = pres.slides.add_slide(pres.slide_layouts[6])
    _add_banner(slide1, base_title + "  ·  产品需求（基础信息）", subtitle_common)
    leaf1 = [
        "序号", "需求编号", "需求标题", "计划版本", "优先级", "所属特性",
        "特性FO", "特性SE", "特性TFO", "涉及代码领域", "关键风险",
    ]
    parents1: List[Optional[tuple]] = [None] * 11
    widths1 = [0.4, 1.0, 1.9, 0.9, 0.55, 1.0, 0.7, 0.7, 0.7, 1.6, 2.5]
    _add_grouped_table(slide1, parents1, leaf1, rows_basic, widths1,
                       center_cols={0, 4, 6, 7, 8})

    # —— Slide 2：交付进展跟踪 ——
    slide2 = pres.slides.add_slide(pres.slide_layouts[6])
    _add_banner(slide2, base_title + "  ·  产品需求（交付进展跟踪）", subtitle_common)
    leaf2 = [
        "序号", "需求编号", "需求标题",
        "需求串讲", "反串讲", "领域串讲", "编码", "联调验证", "转测澄清", "测试结论",
        "预估代码量", "实际代码量", "实际工作量",
    ]
    parents2: List[Optional[tuple]] = [None, None, None] + [("progress", "交付进展跟踪")] * 10
    widths2 = [0.4, 1.0, 1.5, 0.8, 0.7, 0.85, 0.6, 0.85, 0.85, 0.85, 0.95, 0.95, 0.95]
    _add_grouped_table(slide2, parents2, leaf2, rows_progress, widths2,
                       center_cols={0, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12})


def build_iteration_pptx(iteration, requirements: Iterable, product_reqs: Optional[Iterable] = None) -> io.BytesIO:
    """传入 AnnualIteration + 领域/产品 需求列表，返回 BytesIO。

    - 始终输出领域需求 slide（保持原有行为）
    - 若 product_reqs 非空，再追加 2 张产品需求 slide
    """
    pres = _new_pres()
    _add_domain_slide(pres, iteration, requirements)

    if product_reqs:
        product_list = list(product_reqs)
        if product_list:
            _add_product_slides(pres, iteration, product_list)

    buf = io.BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf
