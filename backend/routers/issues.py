"""问题单报表接口：读取服务器侧 Excel，解析后返回给前端。

权限：
- GET  /api/issues/data              —— 所有登录用户（加载最新单张报表）
- GET  /api/issues/trend             —— 所有登录用户（扫描全目录按天聚合趋势）
- GET  /api/issues/run-script/status —— 所有登录用户（查询脚本是否正在运行）
- POST /api/issues/run-script        —— 仅管理员（执行外部刷新脚本）
- GET  /api/issues/export.pptx       —— 所有登录用户（导出 PPT）

配置项（via PUT /api/config）：
  issue_report_path      —— Excel 目录或文件路径
  issue_script_path      —— 刷新脚本路径（.py / .bat / .exe）
  issue_script_timeout   —— 采集脚本超时秒数（默认 600）
  issue_snapshot_time    —— 每日自动采集时刻 HH:MM（默认 07:30）
  issue_snapshot_enabled —— 是否启用每日自动采集（默认 true）
"""
import io
import json
import pathlib
import re
import subprocess
import sys
import threading
from datetime import datetime
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session

import models
from auth import get_current_user, require_admin
from database import SessionLocal, get_db
from op_log import log_op
from routers.config import _load as _load_config

router = APIRouter(prefix="/api/issues", tags=["issues"])

# 全局锁：防止脚本并发执行
_script_lock: threading.Lock = threading.Lock()
_script_started_at: Optional[datetime] = None

_RAW_COLS = [
    "version",          # A 版本信息
    "issue_id",         # B 缺陷业务编号
    "title",            # C 标题
    "owner",            # D 当前责任人
    "group",            # E 当前责任人所属小组
    "progress",         # F 进展
    "severity",         # G 严重程度
    "severity_di",      # H 严重程度DI值
    "root_cause",       # I 根因
    "solution",         # J 解决措施
    "progress_record",  # K 进展记录
    "estimated_close",  # L 预计闭环时间
    "priority",         # M 优先级
    "is_sdts",          # N 是否SDTS
    "year",             # O 年份
    "month",            # P 月份
    "date",             # Q 日期
    "year_month",       # R 年月（钻取按月度过滤的关键字段）
    "category",         # S 标题分类（钻取按客户/分类过滤的关键字段）
    "customer",         # T 客户面（API 快照聚合的关键维度；由后端从标题匹配客户主数据得到）
    "department",       # U 责任人部门（展示用：责任人所在的直属部门）
    "feature",          # V 特性
    "subsystem",        # W 子系统
    "module",           # X 模块
    "dept_path",        # Y 责任人部门全路径（各级部门拼接，仅用于部门过滤匹配，不展示）
]

_DATE_PAT      = re.compile(r"_(\d{8})\.",           re.IGNORECASE)
_DATE_DIR_PAT  = re.compile(r"^\d{4}-\d{2}-\d{2}$")

COLORS = ["#4073ba", "#67C23A", "#E6A23C", "#F56C6C", "#909399", "#8E7AD8", "#26C9C3"]


# ─── helpers ────────────────────────────────────────────────────────────────

def _cell(ws, row: int, col: int) -> str:
    v = ws.cell(row, col).value
    return str(v).strip() if v is not None else ""


def _count_by(rows: List[Dict], field: str) -> Dict[str, int]:
    result: Dict[str, int] = {}
    for r in rows:
        k = r.get(field, "")
        if k:
            result[k] = result.get(k, 0) + 1
    return result


def _file_sort_key(f: pathlib.Path):
    m = _DATE_PAT.search(f.name)
    if m:
        try:
            return (1, datetime.strptime(m.group(1), "%Y%m%d").timestamp())
        except ValueError:
            pass
    return (0, f.stat().st_mtime)


def _resolve_target(path_str: str) -> pathlib.Path:
    p = pathlib.Path(path_str)
    if not p.exists():
        raise HTTPException(404, f"路径不存在：{path_str}")
    if p.is_file():
        if p.suffix.lower() != ".xlsx":
            raise HTTPException(400, "指定文件不是 .xlsx 格式")
        return p
    if p.is_dir():
        candidates = sorted(p.glob("*.xlsx"), key=_file_sort_key, reverse=True)
        if not candidates:
            raise HTTPException(404, f"目录 {path_str} 中未找到 .xlsx 文件")
        return candidates[0]
    raise HTTPException(400, f"无法识别路径：{path_str}")


def _parse_cross_table(ws) -> Dict[str, Any]:
    max_col = ws.max_column
    columns = [_cell(ws, 1, c) for c in range(2, max_col + 1)]
    while columns and not columns[-1]:
        columns.pop()
    rows: List[Dict] = []
    for r in range(2, ws.max_row + 1):
        label = _cell(ws, r, 1)
        if not label:
            break
        row: Dict[str, Any] = {"label": label}
        for i, col in enumerate(columns):
            raw_val = ws.cell(r, i + 2).value
            if raw_val is None:
                row[col] = 0
            else:
                try:
                    row[col] = int(raw_val)
                except (ValueError, TypeError):
                    row[col] = str(raw_val)
        rows.append(row)
    return {"columns": columns, "rows": rows}


def _parse_excel(path: str) -> Dict[str, Any]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
    except FileNotFoundError:
        raise HTTPException(404, "报表文件不存在，请检查路径配置")
    except Exception as exc:
        raise HTTPException(500, f"读取报表文件失败: {exc}")

    raw: List[Dict] = []
    try:
        ws_raw = wb["原始数据"]
        for r in range(2, ws_raw.max_row + 1):
            row = {col: _cell(ws_raw, r, i + 1) for i, col in enumerate(_RAW_COLS)}
            if any(row.values()):
                raw.append(row)
    except KeyError:
        pass

    def _sheet(name):
        try:
            return _parse_cross_table(wb[name])
        except KeyError:
            return {"columns": [], "rows": []}

    try:
        mtime = datetime.fromtimestamp(
            pathlib.Path(path).stat().st_mtime
        ).strftime("%Y-%m-%d %H:%M:%S")
    except Exception:
        mtime = None

    return {
        "file_mtime": mtime,
        "raw": raw,
        "monthly_by_group":   _sheet("按小组月度统计"),
        "annual_by_group":    _sheet("按小组年度统计"),
        "by_customer":        _sheet("按客户统计"),
        "feature_by_group":   _sheet("特性×小组统计"),
        "feature_by_customer":_sheet("特性×客户统计"),
    }


def _parse_raw_from_wb(wb) -> List[Dict]:
    raw: List[Dict] = []
    try:
        ws = wb["原始数据"]
        for r in range(2, ws.max_row + 1):
            row = {col: _cell(ws, r, i + 1) for i, col in enumerate(_RAW_COLS)}
            if any(row.values()):
                raw.append(row)
    except KeyError:
        pass
    return raw


def _list_date_dirs(root: pathlib.Path) -> list:
    """Return date subdirs (YYYY-MM-DD) containing xlsx files, sorted newest-first."""
    result = [
        d for d in root.iterdir()
        if d.is_dir() and _DATE_DIR_PAT.match(d.name) and any(d.glob("*.xlsx"))
    ]
    return sorted(result, key=lambda d: d.name, reverse=True)


def _resolve_for_date(path_str: str, date: Optional[str] = None):
    """Resolve to a single xlsx path, supporting both flat dirs and date-subdir structures."""
    p = pathlib.Path(path_str)
    if not p.exists():
        raise HTTPException(404, f"路径不存在：{path_str}")
    if p.is_file():
        return p

    if p.is_dir():
        date_dirs = _list_date_dirs(p)
        if date_dirs:
            if date:
                target_dir = p / date
                if not target_dir.is_dir():
                    raise HTTPException(404, f"日期目录不存在：{date}")
                xlsxes = sorted(target_dir.glob("*.xlsx"), key=lambda f: f.name, reverse=True)
                if not xlsxes:
                    raise HTTPException(404, f"日期 {date} 目录中无 xlsx 文件")
            else:
                target_dir = date_dirs[0]
                xlsxes = sorted(target_dir.glob("*.xlsx"), key=lambda f: f.name, reverse=True)
            return xlsxes[0]

        # Flat directory fallback
        candidates = sorted(p.glob("*.xlsx"), key=_file_sort_key, reverse=True)
        if not candidates:
            raise HTTPException(404, f"目录 {path_str} 中未找到 .xlsx 文件")
        return candidates[0]

    raise HTTPException(400, f"无法识别路径：{path_str}")


# ─── PPT builder ────────────────────────────────────────────────────────────

def _build_pptx(data: dict) -> io.BytesIO:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    C_BLUE   = RGBColor(0x40, 0x73, 0xBA)
    C_RED    = RGBColor(0xF5, 0x6C, 0x6C)
    C_ORANGE = RGBColor(0xE6, 0xA2, 0x3C)
    C_GRAY   = RGBColor(0x90, 0x93, 0x99)
    C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    C_DARK   = RGBColor(0x30, 0x31, 0x33)
    C_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)

    raw = data.get("raw", [])

    def _txt(slide, text, x, y, w, h, size=12, bold=False, color=C_DARK, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _cell_run(cell, text: str, size: int, bold: bool, color: RGBColor,
                  align=PP_ALIGN.LEFT, bg: RGBColor = None):
        """Set cell text via a run (avoids paragraph.font issues)."""
        if bg is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _table_slide(title: str, columns: List[str], rows_data: List[List[str]]):
        slide = prs.slides.add_slide(blank)
        _txt(slide, title, 0.4, 0.15, 12.5, 0.6, size=20, bold=True, color=C_BLUE)

        n_rows = len(rows_data) + 1
        n_cols = len(columns)
        cell_h = min(0.38, 5.8 / n_rows)
        total_h = cell_h * n_rows
        tbl = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.4), Inches(0.9), Inches(12.5), Inches(total_h)
        ).table

        col_w = [2.0] + [10.5 / max(n_cols - 1, 1)] * (n_cols - 1)

        for c, h in enumerate(columns):
            _cell_run(tbl.cell(0, c), h, 10, True, C_WHITE,
                      align=PP_ALIGN.CENTER, bg=C_BLUE)

        for r, row_vals in enumerate(rows_data):
            is_total = bool(row_vals) and row_vals[0] == "合计"
            for c, val in enumerate(row_vals):
                _cell_run(tbl.cell(r + 1, c), str(val), 10, is_total, C_DARK,
                          align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT,
                          bg=C_LIGHT if is_total else None)

        for c, w in enumerate(col_w):
            tbl.columns[c].width = Inches(w)

    # ── Slide 1: Title ────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank)
    _txt(slide1, "缺陷统计报表", 1, 1.8, 11, 1.4, size=48, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    _txt(slide1, f"{data.get('actual_file', '')}   {data.get('file_mtime', '')}",
         1, 3.4, 11, 0.5, size=13, color=C_GRAY, align=PP_ALIGN.CENTER)

    total  = len(raw)
    cnts   = {s: sum(1 for r in raw if r.get("severity") == s) for s in ["严重", "一般", "提示"]}
    cards  = [("合计", total, C_BLUE), ("严重", cnts["严重"], C_RED),
              ("一般", cnts["一般"], C_ORANGE), ("提示", cnts["提示"], C_GRAY)]
    cx = 0.7
    for label, val, clr in cards:
        shp = slide1.shapes.add_shape(1, Inches(cx), Inches(4.4), Inches(2.8), Inches(1.7))
        shp.fill.solid()
        shp.fill.fore_color.rgb = clr
        shp.line.color.rgb = clr
        tf = shp.text_frame
        tf.word_wrap = False
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = str(val)
        r1.font.size = Pt(40)
        r1.font.bold = True
        r1.font.color.rgb = C_WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(14)
        r2.font.color.rgb = C_WHITE
        cx += 3.0

    # ── Slide 2: Monthly by group ─────────────────────────────────
    m = data.get("monthly_by_group", {})
    if m.get("rows"):
        cols  = ["小组"] + m["columns"]
        rdata = [[r["label"]] + [str(r.get(c, 0)) for c in m["columns"]] for r in m["rows"]]
        _table_slide("按小组月度统计", cols, rdata)

    # ── Slide 3: By customer ─────────────────────────────────────
    c = data.get("by_customer", {})
    if c.get("rows"):
        cols  = ["小组"] + c["columns"]
        rdata = [[r["label"]] + [str(r.get(col, 0)) for col in c["columns"]] for r in c["rows"]]
        _table_slide("按客户分布", cols, rdata)

    # ── Slide 4: Feature by group ─────────────────────────────────
    f = data.get("feature_by_group", {})
    if f.get("rows"):
        cols  = ["小组"] + f["columns"]
        rdata = [[r["label"]] + [str(r.get(col, 0)) for col in f["columns"]] for r in f["rows"]]
        _table_slide("特性 × 小组分布", cols, rdata)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─── endpoints ──────────────────────────────────────────────────────────────

@router.get("/dates")
def list_dates(_: models.User = Depends(get_current_user)):
    """列出报表目录下所有含 xlsx 文件的日期子目录（YYYY-MM-DD），最新在前。"""
    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        return []
    p = pathlib.Path(path_str)
    if not p.is_dir():
        return []
    return [d.name for d in _list_date_dirs(p)]


@router.get("/data")
def get_data(date: Optional[str] = None, _: models.User = Depends(get_current_user)):
    """加载单张报表。date 参数为 YYYY-MM-DD，不传则取最新一天。"""
    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        return {"configured": False}
    target = _resolve_for_date(path_str, date)
    result = _parse_excel(str(target))
    result["actual_file"] = target.name
    result["date_dir"] = target.parent.name if _DATE_DIR_PAT.match(target.parent.name) else None
    return {"configured": True, **result}


@router.get("/trend")
def get_trend(_: models.User = Depends(get_current_user)):
    """扫描目录内全部报表，按天聚合趋势数据。支持日期子目录和平铺两种结构。"""
    import openpyxl

    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        raise HTTPException(400, "未配置报表路径")

    p = pathlib.Path(path_str)
    if p.is_file():
        p = p.parent
    if not p.is_dir():
        raise HTTPException(404, f"目录不存在：{path_str}")

    # 日期子目录结构（优先）
    date_dirs = _list_date_dirs(p)
    if date_dirs:
        file_list = []
        for d in reversed(date_dirs):  # 升序用于趋势
            xlsxes = sorted(d.glob("*.xlsx"), key=lambda f: f.name, reverse=True)
            if xlsxes:
                file_list.append((d.name, xlsxes[0]))
        if not file_list:
            raise HTTPException(404, "无有效报表数据")
    else:
        # 平铺目录兼容
        flat = sorted(
            (f for f in p.glob("*.xlsx") if _DATE_PAT.search(f.name)),
            key=_file_sort_key,
        )
        if not flat:
            raise HTTPException(404, "目录中无含日期后缀的报表文件")
        file_list = []
        for f in flat:
            m = _DATE_PAT.search(f.name)
            ds = m.group(1)
            file_list.append((f"{ds[:4]}-{ds[4:6]}-{ds[6:]}", f))

    daily: List[Dict] = []
    all_groups:     set = set()
    all_severities: set = set()

    for date_str, f in file_list:
        try:
            wb  = openpyxl.load_workbook(str(f), data_only=True)
            raw = _parse_raw_from_wb(wb)
            wb.close()
        except Exception:
            continue

        bg = _count_by(raw, "group")
        bs = _count_by(raw, "severity")
        all_groups.update(bg.keys())
        all_severities.update(bs.keys())

        daily.append({
            "date":        date_str,
            "file":        f.name,
            "total":       len(raw),
            "by_group":    bg,
            "by_severity": bs,
        })

    sev_order = ["严重", "一般", "提示"]
    return {
        "daily":          daily,
        "all_groups":     sorted(all_groups),
        "all_severities": sorted(all_severities, key=lambda s: sev_order.index(s) if s in sev_order else 99),
    }


# ─── 通过脚本调用外部 API 拉取问题单（按项目）──────────────────────────────────
def _normalize_issue_row(r: dict) -> Dict[str, str]:
    """把脚本返回的一条问题单规整成「原始数据」表同款字段（缺的留空）。"""
    return {col: (str(r.get(col)).strip() if r.get(col) is not None else "") for col in _RAW_COLS}


def _script_timeout(cfg: Dict) -> int:
    """脚本执行超时（秒）。config.issue_script_timeout，默认 600s。

    原来写死 120s：DTS 接口慢一点就被杀，但快照有时已经落盘，表现为"报错了、
    过会儿刷新数据又有了"。留成可配置，默认给足。
    """
    try:
        v = int(cfg.get("issue_script_timeout") or 0)
    except (TypeError, ValueError):
        v = 0
    return v if v > 0 else 600


def _run_issue_api_script(project: str) -> List[Dict]:
    """以 `python <issue_api_script_path> <project>` 调用脚本，期望 stdout 为 JSON 数组。"""
    cfg = _load_config()
    script = (cfg.get("issue_api_script_path") or "").strip()
    if not script:
        raise HTTPException(400, "未配置 API 脚本（issue_api_script_path）")
    sp = pathlib.Path(script)
    if not sp.exists():
        raise HTTPException(404, f"脚本不存在：{script}")

    timeout = _script_timeout(cfg)
    cmd = [sys.executable, str(sp), project] if sp.suffix.lower() == ".py" else [str(sp), project]
    try:
        result = subprocess.run(
            cmd, capture_output=True, encoding="utf-8", errors="replace",
            timeout=timeout, cwd=str(sp.parent),
        )
    except subprocess.TimeoutExpired:
        raise HTTPException(504, f"脚本执行超时（>{timeout} 秒），可在「配置」中调大超时时间")
    except Exception as exc:
        raise HTTPException(500, f"脚本启动失败：{exc}")

    if result.returncode != 0:
        raise HTTPException(500, f"脚本退出码 {result.returncode}：{(result.stderr or '')[-500:]}")
    out = (result.stdout or "").strip()
    if not out:
        raise HTTPException(500, "脚本无输出（应向 stdout 打印 JSON 数组）")
    try:
        data = json.loads(out)
    except json.JSONDecodeError as exc:
        raise HTTPException(500, f"脚本输出不是合法 JSON：{exc}")
    if not isinstance(data, list):
        raise HTTPException(500, "脚本输出应为问题单数组（JSON list）")
    return [_normalize_issue_row(r) for r in data if isinstance(r, dict)]


@router.get("/api-data")
def get_api_data(project: str, _: models.User = Depends(get_current_user)):
    """按项目（YLS3000 / YLS5000 / YLS8000）通过脚本调用外部 API 拉取问题单。

    脚本契约：后端以 `python <issue_api_script_path> <project>` 调用，脚本把问题单
    列表以 JSON 数组打印到 stdout（字段同「原始数据」表：version/issue_id/title/owner/
    group/progress/severity/…）。失败时以 200 + error 字段返回，便于前端友好提示。
    """
    cfg = _load_config()
    if not (cfg.get("issue_api_script_path") or "").strip():
        return {"configured": False, "project": project}
    try:
        raw = _run_issue_api_script(project)
    except HTTPException as exc:
        return {
            "configured": True, "project": project, "error": str(exc.detail),
            "count": 0, "raw": [], "by_severity": {}, "by_group": {}, "by_customer": {},
        }
    return {
        "configured": True,
        "project": project,
        "fetched_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "count": len(raw),
        "raw": raw,
        "by_severity": _count_by(raw, "severity"),
        "by_group": _count_by(raw, "group"),
        "by_customer": _count_by(raw, "category"),
        "error": None,
    }


# ─── 每日快照：库存"数字"（趋势）+ 文件存明细（钻取）──────────────────────────
_BACKEND_DIR = pathlib.Path(__file__).resolve().parent.parent
_SEV_ORDER = {"严重": 0, "一般": 1, "提示": 2}


def _snapshot_root() -> pathlib.Path:
    """快照明细文件根目录：优先 config.issue_snapshot_dir，否则 backend/data/issue_snapshots。"""
    cfg = _load_config()
    d = (cfg.get("issue_snapshot_dir") or "").strip()
    root = pathlib.Path(d) if d else (_BACKEND_DIR / "data" / "issue_snapshots")
    root.mkdir(parents=True, exist_ok=True)
    return root


def _safe_slug(s: str) -> str:
    return re.sub(r"[^\w\-]", "_", s or "") or "_"


# ─── 采集后富化：部门过滤 + 责任人归组 + 从标题提取客户面 ─────────────────────
def _as_str_list(v) -> List[str]:
    """config 里可能存成 list 或分号/换行分隔的字符串，统一成去空的列表。"""
    if isinstance(v, list):
        return [str(x).strip() for x in v if str(x).strip()]
    if isinstance(v, str):
        return [s.strip() for s in re.split(r"[;；\n]", v) if s.strip()]
    return []


def _load_issue_groups(cfg: Dict) -> List[tuple]:
    """config.issue_groups: [{name, members}] → [(小组名, [成员,...])]。成员支持分号/换行分隔。"""
    groups: List[tuple] = []
    for g in (cfg.get("issue_groups") or []):
        if not isinstance(g, dict):
            continue
        name = str(g.get("name") or "").strip()
        if not name:
            continue
        groups.append((name, _as_str_list(g.get("members"))))
    return groups


def _match_group(owner: str, groups: List[tuple]) -> str:
    """按责任人姓名匹配所属小组（大小写不敏感 + 互为子串的模糊匹配）。"""
    o = (owner or "").strip().lower()
    if not o:
        return ""
    for name, members in groups:
        for m in members:
            ml = m.lower()
            if ml and (ml == o or ml in o or o in ml):
                return name
    return ""


def _load_customer_matchers(db: Session) -> List[tuple]:
    """从客户主数据（code/全称/别名）构建 [(匹配文本_lower, 展示名)]，按长度降序（优先更具体）。"""
    matchers: List[tuple] = []
    customers = db.query(models.Customer).filter(models.Customer.is_active == True).all()  # noqa: E712
    id2label = {}
    for c in customers:
        label = (c.display_name or c.code or "").strip()
        id2label[c.id] = label
        for t in (c.code, c.display_name):
            if t and t.strip():
                matchers.append((t.strip().lower(), label))
    for a in db.query(models.CustomerAlias).all():
        label = id2label.get(a.customer_id)
        if label and a.alias and a.alias.strip():
            matchers.append((a.alias.strip().lower(), label))
    seen, uniq = set(), []
    for mt, label in matchers:
        if mt and mt not in seen:
            seen.add(mt)
            uniq.append((mt, label))
    uniq.sort(key=lambda x: len(x[0]), reverse=True)
    return uniq


def _match_customer(title: str, matchers: List[tuple]) -> str:
    t = (title or "").lower()
    if not t:
        return ""
    for mt, label in matchers:
        if mt in t:
            return label
    return ""


def _enrich_rows(db: Session, rows: List[Dict]) -> List[Dict]:
    """对采集到的问题单做：① 部门过滤 ② 按责任人归组 ③ 从标题提取客户面。

    配置项（config.json，问题单管理「配置」tab 维护）：
      issue_exclude_statuses —— 直接剔除的状态（默认 关闭/撤销；子串匹配进展/状态）
      issue_stat_departments —— 只统计这些部门（子串匹配责任人部门全路径；留空＝全部）
      issue_groups           —— [{name, members}]，成员分号分隔，按责任人归组
    客户面来自客户主数据（客户面管理），用 code/全称/别名 在标题里做包含匹配。
    """
    cfg = _load_config()
    exclude = _as_str_list(cfg.get("issue_exclude_statuses")) or ["关闭", "撤销"]
    depts = _as_str_list(cfg.get("issue_stat_departments"))
    groups = _load_issue_groups(cfg)
    matchers = _load_customer_matchers(db)

    out: List[Dict] = []
    for r in rows:
        # ① 先剔除已关闭 / 已撤销 的单（任何地方都不出现）
        prog = r.get("progress", "") or ""
        if any(s and s in prog for s in exclude):
            continue
        # ② 部门过滤：匹配责任人部门全路径（兼容部门落在上级字段的情况），回退直属部门
        if depts:
            dept = (r.get("dept_path") or "") or (r.get("department") or "")
            if not any(d in dept for d in depts):
                continue
        # ③ 按责任人归组（不在任何小组名单的不保留）
        if groups:
            g = _match_group(r.get("owner", ""), groups)
            if not g:
                continue
            r["group"] = g
        # ④ 从标题提取客户面
        if matchers and not (r.get("customer") or "").strip():
            r["customer"] = _match_customer(r.get("title", ""), matchers)
        out.append(r)
    return out


def _take_snapshot(db: Session, project: str, source: str = "api") -> models.IssueSnapshot:
    """拉取该项目问题单 → 明细写文件、聚合数字写库（同项目同日覆盖）。

    可能抛 HTTPException（脚本未配置 / 执行失败）——调用方按需捕获。
    """
    raw = _run_issue_api_script(project)
    raw = _enrich_rows(db, raw)   # 部门过滤 + 责任人归组 + 标题提取客户面
    today = datetime.now().strftime("%Y-%m-%d")

    # 明细落文件（<项目>/<日期>.json）
    root = _snapshot_root()
    rel = f"{_safe_slug(project)}/{today}.json"
    fp = root / rel
    fp.parent.mkdir(parents=True, exist_ok=True)
    fp.write_text(json.dumps(raw, ensure_ascii=False), encoding="utf-8")

    # 自动落盘 Excel 备份：原始表 + 分析表（同日多次用时间戳不覆盖）
    _export_snapshot_excel(project, raw, today)

    # upsert 快照元数据
    snap = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project,
                models.IssueSnapshot.snapshot_date == today)
        .first()
    )
    if snap is None:
        snap = models.IssueSnapshot(project=project, snapshot_date=today)
        db.add(snap)
    snap.total = len(raw)
    snap.data_file = rel
    snap.source = source
    snap.created_at = datetime.utcnow()
    db.flush()  # 拿到 snap.id

    # 重建维度聚合数字（group / customer / severity）
    db.query(models.IssueSnapshotStat).filter(
        models.IssueSnapshotStat.snapshot_id == snap.id
    ).delete(synchronize_session=False)
    for dim in ("group", "customer", "severity"):
        for key, cnt in _count_by(raw, dim).items():
            db.add(models.IssueSnapshotStat(
                snapshot_id=snap.id, dimension=dim, dim_key=key, count=cnt,
            ))
    db.commit()
    db.refresh(snap)
    return snap


def collect_with_log(db: Session, project: str, source: str = "auto") -> Dict[str, Any]:
    """采集一个项目并写执行日志（成功失败都写）。不抛异常，返回结果字典。

    定时任务与手动采集共用，保证两条路径的日志口径一致。
    """
    started = datetime.utcnow()
    log = models.IssueCollectLog(project=project, source=source, started_at=started)
    result: Dict[str, Any]
    try:
        snap = _take_snapshot(db, project, source=source)
        log.ok, log.total, log.error = True, snap.total, ""
        result = {"project": project, "ok": True, "date": snap.snapshot_date, "total": snap.total}
    except HTTPException as exc:
        log.ok, log.error = False, str(exc.detail)[:2000]
        result = {"project": project, "ok": False, "error": log.error}
    except Exception as exc:  # noqa: BLE001 —— 脚本/解析的意外错误也要留痕
        log.ok, log.error = False, f"{type(exc).__name__}: {exc}"[:2000]
        result = {"project": project, "ok": False, "error": log.error}

    log.finished_at = datetime.utcnow()
    log.duration_ms = int((log.finished_at - started).total_seconds() * 1000)
    try:
        db.add(log)
        db.commit()
    except Exception:  # 日志写失败不能影响采集结论
        db.rollback()
    return result


# ─── 手动采集：后台线程执行 + 轮询状态 ────────────────────────────────────────
# 采集要跑几分钟，而前端 axios 超时只有 10s：同步返回必然先弹「采集失败」、
# 但后台其实还在跑并最终成功，于是"过会儿刷新数据又有了"。改成立即返回 + 轮询。
_collect_lock = threading.Lock()
_collect_state: Dict[str, Any] = {
    "running": False, "projects": [], "current": None,
    "started_at": None, "finished_at": None, "results": [],
}


def _collect_worker(projects: List[str], source: str) -> None:
    db = SessionLocal()
    try:
        for p in projects:
            _collect_state["current"] = p
            _collect_state["results"].append(collect_with_log(db, p, source=source))
    finally:
        db.close()
        _collect_state["current"] = None
        _collect_state["finished_at"] = datetime.utcnow().isoformat()
        _collect_state["running"] = False
        if _collect_lock.locked():
            _collect_lock.release()


@router.post("/snapshot-collect")
def snapshot_collect(
    request: Request,
    project: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(require_admin),
):
    """启动一次采集（仅管理员），立即返回；用 GET /collect-status 轮询进度。

    project 省略则采集 config.issue_api_projects 全部。
    """
    cfg = _load_config()
    projects = [project] if project else (cfg.get("issue_api_projects") or [])
    if not projects:
        raise HTTPException(400, "没有可采集的项目（未配置 issue_api_projects）")

    if not _collect_lock.acquire(blocking=False):
        raise HTTPException(423, f"已有采集任务在执行（{_collect_state.get('current') or '…'}），请稍候")

    _collect_state.update({
        "running": True, "projects": list(projects), "current": None,
        "started_at": datetime.utcnow().isoformat(), "finished_at": None, "results": [],
    })
    threading.Thread(target=_collect_worker, args=(list(projects), "manual"), daemon=True).start()

    log_op(db, action="issue_snapshot", target="issue_snapshot", target_id=None,
           detail=f"手动触发采集：{', '.join(projects)}", user=current_user, request=request)
    return {"started": True, "projects": projects}


@router.get("/collect-status")
def collect_status(_: models.User = Depends(get_current_user)):
    """采集任务进度（所有登录用户可查）。running=false 且 results 非空即为本轮结果。"""
    return dict(_collect_state)


@router.get("/collect-logs")
def collect_logs(project: Optional[str] = None, limit: int = 50,
                 db: Session = Depends(get_db),
                 _: models.User = Depends(get_current_user)):
    """采集执行日志（新→旧）。project 省略则返回所有项目。"""
    q = db.query(models.IssueCollectLog)
    if project:
        q = q.filter(models.IssueCollectLog.project == project)
    rows = q.order_by(models.IssueCollectLog.started_at.desc()).limit(max(1, min(limit, 200))).all()
    return [
        {
            "id": r.id, "project": r.project, "source": r.source, "ok": bool(r.ok),
            "total": r.total or 0, "duration_ms": r.duration_ms or 0, "error": r.error or "",
            "started_at": r.started_at.strftime("%Y-%m-%d %H:%M:%S") if r.started_at else "",
            "finished_at": r.finished_at.strftime("%Y-%m-%d %H:%M:%S") if r.finished_at else "",
        }
        for r in rows
    ]


@router.get("/snapshots")
def list_snapshots(project: str, db: Session = Depends(get_db),
                   _: models.User = Depends(get_current_user)):
    """某项目的快照列表（新→旧），只含元数据数字。"""
    rows = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project)
        .order_by(models.IssueSnapshot.snapshot_date.desc())
        .all()
    )
    return [
        {"id": r.id, "date": r.snapshot_date, "total": r.total, "source": r.source,
         "created_at": r.created_at.strftime("%Y-%m-%d %H:%M:%S") if r.created_at else ""}
        for r in rows
    ]


@router.get("/snapshot-detail")
def snapshot_detail(project: str, date: Optional[str] = None,
                    db: Session = Depends(get_db),
                    _: models.User = Depends(get_current_user)):
    """某次快照的明细（从文件加载完整行）；date 省略取最新。"""
    q = db.query(models.IssueSnapshot).filter(models.IssueSnapshot.project == project)
    snap = (q.filter(models.IssueSnapshot.snapshot_date == date).first() if date
            else q.order_by(models.IssueSnapshot.snapshot_date.desc()).first())
    if snap is None:
        return {"exists": False, "project": project}
    raw: List[Dict] = []
    try:
        fp = _snapshot_root() / snap.data_file
        if fp.exists():
            raw = json.loads(fp.read_text(encoding="utf-8"))
    except Exception:
        raw = []
    return {
        "exists": True, "project": project, "date": snap.snapshot_date,
        "created_at": snap.created_at.strftime("%Y-%m-%d %H:%M:%S") if snap.created_at else "",
        "total": snap.total, "count": len(raw), "raw": raw,
        "by_severity": _count_by(raw, "severity"),
        "by_group": _count_by(raw, "group"),
        "by_customer": _count_by(raw, "customer"),
    }


@router.get("/snapshot-trend")
def snapshot_trend(project: str, dimension: str = "group",
                   db: Session = Depends(get_db),
                   _: models.User = Depends(get_current_user)):
    """趋势：只从库里读维度聚合数字（不碰明细文件）。dimension ∈ group/customer/severity。"""
    if dimension not in ("group", "customer", "severity"):
        dimension = "group"
    snaps = (
        db.query(models.IssueSnapshot)
        .filter(models.IssueSnapshot.project == project)
        .order_by(models.IssueSnapshot.snapshot_date.asc())
        .all()
    )
    if not snaps:
        return {"project": project, "dimension": dimension, "dates": [], "total": [], "series": []}

    dates = [s.snapshot_date for s in snaps]
    total = [s.total for s in snaps]
    id_to_idx = {s.id: i for i, s in enumerate(snaps)}
    stats = (
        db.query(models.IssueSnapshotStat)
        .filter(models.IssueSnapshotStat.dimension == dimension,
                models.IssueSnapshotStat.snapshot_id.in_([s.id for s in snaps]))
        .all()
    )
    matrix: Dict[str, List[int]] = {}
    order: List[str] = []
    for st in stats:
        idx = id_to_idx.get(st.snapshot_id)
        if idx is None:
            continue
        if st.dim_key not in matrix:
            matrix[st.dim_key] = [0] * len(dates)
            order.append(st.dim_key)
        matrix[st.dim_key][idx] = st.count
    if dimension == "severity":
        order.sort(key=lambda k: _SEV_ORDER.get(k, 99))
    else:
        order.sort()
    series = [{"name": k, "data": matrix[k]} for k in order]
    return {"project": project, "dimension": dimension, "dates": dates,
            "total": total, "series": series}


# ─── 快照导出 Excel：原始数据 + 统计分析 两张表 ────────────────────────────────
def _cross_table(rows: List[Dict], row_field: str, col_field: str,
                 col_order: Optional[List[str]] = None,
                 row_fallback: str = "未标注") -> Dict[str, Any]:
    """行维度 × 列维度 交叉计数 → {columns:[...,'合计'], rows:[{label,...}], total_row}。"""
    matrix: Dict[str, Dict[str, int]] = {}
    col_totals: Dict[str, int] = {}
    col_seen: List[str] = []
    grand = 0
    for r in rows:
        rv = (r.get(row_field) or "").strip() or row_fallback
        cv = (r.get(col_field) or "").strip() or "未标注"
        if cv not in col_totals:
            col_totals[cv] = 0
            col_seen.append(cv)
        matrix.setdefault(rv, {})
        matrix[rv][cv] = matrix[rv].get(cv, 0) + 1
        col_totals[cv] += 1
        grand += 1
    if col_order:
        cols = [c for c in col_order if c in col_totals] + sorted(c for c in col_seen if c not in col_order)
    else:
        cols = sorted(col_seen)
    columns = cols + ["合计"]
    out_rows = []
    for rv in sorted(matrix.keys()):
        rec: Dict[str, Any] = {"label": rv}
        t = 0
        for c in cols:
            rec[c] = matrix[rv].get(c, 0)
            t += rec[c]
        rec["合计"] = t
        out_rows.append(rec)
    total_row = {"label": "合计", **{c: col_totals.get(c, 0) for c in cols}, "合计": grand}
    return {"columns": columns, "rows": out_rows, "total_row": total_row}


_RAW_XLSX_COLS = [
    ("issue_id", "缺陷业务编号", 20), ("title", "标题", 42), ("owner", "当前责任人", 12),
    ("group", "所属小组", 14), ("department", "责任人部门", 20), ("customer", "客户面", 14),
    ("feature", "特性", 14), ("subsystem", "子系统", 14), ("module", "模块", 14),
    ("progress", "进展", 12), ("severity", "严重程度", 10), ("year_month", "年月", 10),
    ("version", "版本信息", 22),
]


def _fill_raw_sheet(ws, raw: List[Dict]) -> None:
    """原始数据表：问题单明细，一行一条。"""
    from openpyxl.styles import Alignment, Font, PatternFill
    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="4073BA")
    center = Alignment(horizontal="center", vertical="center")
    ws.title = "原始数据"
    ws.append([h for _, h, _ in _RAW_XLSX_COLS])
    for c_idx, (_, _, w) in enumerate(_RAW_XLSX_COLS, start=1):
        cell = ws.cell(1, c_idx)
        cell.font = head_font
        cell.fill = head_fill
        cell.alignment = center
        ws.column_dimensions[cell.column_letter].width = w
    for r in raw:
        ws.append([r.get(k, "") for k, _, _ in _RAW_XLSX_COLS])
    ws.freeze_panes = "A2"


def _fill_analysis_sheet(ws, raw: List[Dict]) -> None:
    """统计分析表：按小组 / 客户面 / 年月 × 严重程度 三张交叉表纵向排布。"""
    from openpyxl.styles import Alignment, Font, PatternFill
    head_font = Font(bold=True, color="FFFFFF")
    head_fill = PatternFill("solid", fgColor="4073BA")
    title_font = Font(bold=True, size=12, color="4073BA")
    total_font = Font(bold=True)
    center = Alignment(horizontal="center", vertical="center")
    ws.title = "统计分析"
    ws.column_dimensions["A"].width = 18
    SEV = ["严重", "一般", "提示"]
    row_ptr = [1]   # 显式维护当前写入行，避免依赖 max_row

    def _write_cross(row_label: str, title: str, cross: Dict[str, Any]):
        r = row_ptr[0]
        ws.cell(r, 1, title).font = title_font
        r += 1
        for c_idx, val in enumerate([row_label] + cross["columns"], start=1):
            cell = ws.cell(r, c_idx, val)
            cell.font = head_font
            cell.fill = head_fill
            cell.alignment = center
        for row in cross["rows"]:
            r += 1
            ws.cell(r, 1, row["label"])
            for c_idx, col in enumerate(cross["columns"], start=2):
                ws.cell(r, c_idx, row.get(col, 0)).alignment = center
        r += 1
        ws.cell(r, 1, cross["total_row"]["label"]).font = total_font
        for c_idx, col in enumerate(cross["columns"], start=2):
            cell = ws.cell(r, c_idx, cross["total_row"].get(col, 0))
            cell.font = total_font
            cell.alignment = center
        row_ptr[0] = r + 2   # 空一行再写下一张表

    _write_cross("小组", "按小组 × 严重程度", _cross_table(raw, "group", "severity", SEV, "未分组"))
    _write_cross("客户面", "按客户面 × 严重程度", _cross_table(raw, "customer", "severity", SEV, "未标注"))
    _write_cross("年月", "按年月 × 严重程度", _cross_table(raw, "year_month", "severity", SEV, "未标注"))


def _excel_dir(which: str) -> pathlib.Path:
    """自动导出目录：which=raw（原始表）/ analysis（分析表），各自可配置，
    默认 backend/data/issue_excel/<which>。"""
    cfg = _load_config()
    key = "issue_excel_raw_dir" if which == "raw" else "issue_excel_analysis_dir"
    d = (cfg.get(key) or "").strip()
    root = pathlib.Path(d) if d else (_BACKEND_DIR / "data" / "issue_excel" / which)
    root.mkdir(parents=True, exist_ok=True)
    return root


def _export_snapshot_excel(project: str, raw: List[Dict], date_str: str) -> None:
    """采集后自动落盘：原始表 + 分析表 各存一份到备份目录（同日多次用时间戳，不覆盖）。

    失败只吞掉不影响采集主流程。
    """
    try:
        import openpyxl
        slug = _safe_slug(project)
        stem = f"{date_str}_{datetime.now().strftime('%H%M%S')}"

        raw_wb = openpyxl.Workbook()
        _fill_raw_sheet(raw_wb.active, raw)
        raw_dir = _excel_dir("raw") / slug
        raw_dir.mkdir(parents=True, exist_ok=True)
        raw_wb.save(str(raw_dir / f"{slug}_原始_{stem}.xlsx"))

        ana_wb = openpyxl.Workbook()
        _fill_analysis_sheet(ana_wb.active, raw)
        ana_dir = _excel_dir("analysis") / slug
        ana_dir.mkdir(parents=True, exist_ok=True)
        ana_wb.save(str(ana_dir / f"{slug}_分析_{stem}.xlsx"))
    except Exception:
        pass


@router.get("/snapshot-export")
def snapshot_export(request: Request, project: str, date: Optional[str] = None,
                    db: Session = Depends(get_db),
                    current_user: models.User = Depends(get_current_user)):
    """把某次快照导出为 Excel：Sheet1「原始数据」+ Sheet2「统计分析」（按小组/客户面/年月 × 严重程度）。"""
    import openpyxl

    q = db.query(models.IssueSnapshot).filter(models.IssueSnapshot.project == project)
    snap = (q.filter(models.IssueSnapshot.snapshot_date == date).first() if date
            else q.order_by(models.IssueSnapshot.snapshot_date.desc()).first())
    if snap is None:
        raise HTTPException(404, "该项目暂无快照可导出")
    raw: List[Dict] = []
    try:
        fp = _snapshot_root() / snap.data_file
        if fp.exists():
            raw = json.loads(fp.read_text(encoding="utf-8"))
    except Exception:
        raw = []

    wb = openpyxl.Workbook()
    _fill_raw_sheet(wb.active, raw)
    _fill_analysis_sheet(wb.create_sheet(), raw)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    filename = f"issues_{_safe_slug(project)}_{snap.snapshot_date}.xlsx"
    log_op(db, action="导出Excel", target="问题单", target_id=snap.id,
           detail=f"project={project} date={snap.snapshot_date} rows={len(raw)}",
           user=current_user, request=request)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.get("/run-script/status")
def run_script_status(_: models.User = Depends(get_current_user)):
    """查询刷新脚本是否正在执行（所有登录用户可查）。"""
    return {
        "running":    _script_lock.locked(),
        "started_at": _script_started_at.isoformat() if _script_started_at else None,
    }


@router.post("/run-script")
def run_script(
    request: Request,
    db: Session = Depends(get_db),
    current_admin: models.User = Depends(require_admin),
):
    """执行管理员配置的外部刷新脚本（全局互斥，同时只能有一个实例）。"""
    global _script_started_at

    if not _script_lock.acquire(blocking=False):
        raise HTTPException(423, "脚本正在执行中，请等待完成后再试")

    _script_started_at = datetime.utcnow()
    try:
        cfg = _load_config()
        script = cfg.get("issue_script_path", "").strip()
        if not script:
            raise HTTPException(400, "未配置刷新脚本路径（issue_script_path）")

        sp = pathlib.Path(script)
        if not sp.exists():
            raise HTTPException(404, f"脚本不存在：{script}")

        cmd = [sys.executable, str(sp)] if sp.suffix.lower() == ".py" else [str(sp)]

        result = subprocess.run(
            cmd, capture_output=True, text=True,
            timeout=300, cwd=str(sp.parent),
        )
        log_op(db, action="运行脚本", target="问题单",
               detail=f"script={sp.name} exit={result.returncode}",
               user=current_admin, request=request)
        return {
            "ok":        result.returncode == 0,
            "exit_code": result.returncode,
            "stdout":    result.stdout[-3000:] if result.stdout else "",
            "stderr":    result.stderr[-1000:] if result.stderr else "",
        }
    except HTTPException:
        raise
    except subprocess.TimeoutExpired:
        raise HTTPException(500, "脚本执行超时（>5分钟）")
    except Exception as exc:
        raise HTTPException(500, f"脚本启动失败：{exc}")
    finally:
        _script_started_at = None
        _script_lock.release()


@router.get("/export.pptx")
def export_pptx(
    request: Request,
    date: Optional[str] = None,
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user),
):
    """将指定日期报表导出为 PPT，不传 date 则取最新。"""
    cfg = _load_config()
    path_str = cfg.get("issue_report_path", "").strip()
    if not path_str:
        raise HTTPException(400, "未配置报表路径")

    target = _resolve_for_date(path_str, date)
    data   = _parse_excel(str(target))
    data["actual_file"] = target.name

    try:
        buf = _build_pptx(data)
    except Exception as exc:
        raise HTTPException(500, f"PPT 生成失败：{exc}")
    filename = f"缺陷统计报表_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    log_op(db, action="导出PPT", target="问题单",
           detail=f"date={date or '最新'} file={target.name}",
           user=current_user, request=request)
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
